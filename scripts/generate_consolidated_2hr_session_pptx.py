#!/usr/bin/env python3
"""
Build one consolidated PowerPoint for the 2-hour MLOps session.

Order (per Session Key Points image):
  1) Jenkins in MLOps — CI/CD and automated deployment
  2) Flask ML API for E-Commerce
  3) Cloud-Driven ML Pipeline for Financial Prediction
  4) Automated ML Pipelines and Workflow Automation

Then: hands-on execution slides for the mlops-workshop repo.

Content is synthesized from existing case-study decks in ppts/:
  - case study.pptx (Jenkins / CD in MLOps narrative)
  - case study 3.pptx (Flask e-commerce API)
  - case study (3).pptx (financial / cloud ML pipeline, quantization)
  - case study (2).pptx (AutoML / automated pipelines)
  - mlops-security-log-anomaly-usecase.pptx + data-security-usecase-demo-slides.pptx (security bridge)

Requires: python-pptx
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# --- Theme (clean, presentation-friendly) ---
C_BG = RGBColor(255, 255, 255)
C_NAVY = RGBColor(26, 35, 126)
C_TEAL = RGBColor(0, 121, 107)
C_MUTED = RGBColor(55, 55, 65)
C_CODE_BG = RGBColor(245, 248, 252)

# Fonts: Calibri is the Office default and renders consistently on Windows/macOS.
FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_CODE = "Consolas"


def _blank_slide(prs: Presentation):
    # Layout 6 is typically "Blank" in default Office theme
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _accent_bar(slide, top=Inches(0), height=Inches(0.22)) -> None:
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, prs_w(slide), height
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = C_NAVY
    shp.line.fill.background()


def prs_w(slide) -> int:
    return slide.part.package.presentation_part.presentation.slide_width


def _style_body_paragraph(p, size_pt: int, bold: bool = False) -> None:
    p.font.name = FONT_BODY
    p.font.size = Pt(size_pt)
    p.font.color.rgb = C_MUTED
    p.font.bold = bold
    # Comfortable line height (multiple of font size)
    try:
        p.line_spacing = 1.18
    except (TypeError, AttributeError):
        pass
    p.space_after = Pt(10)
    p.space_before = Pt(0)


def _title_block(
    slide,
    title: str,
    subtitle: str | None = None,
    top=Inches(0.45),
) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), top, Inches(12.0), Inches(1.55))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = C_NAVY
    p.alignment = PP_ALIGN.LEFT
    try:
        p.line_spacing = 1.12
    except (TypeError, AttributeError):
        pass
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT_TITLE
        p2.font.size = Pt(17)
        p2.font.color.rgb = C_MUTED
        p2.font.bold = False
        p2.space_before = Pt(8)
        p2.space_after = Pt(0)
        try:
            p2.line_spacing = 1.2
        except (TypeError, AttributeError):
            pass


def _bullets(
    slide,
    items: list[str],
    top=Inches(1.78),
    left=Inches(0.7),
    width=Inches(12.0),
    height=Inches(5.15),
    size=19,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(6)
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        _style_body_paragraph(p, size, bold=False)


def _command_block(
    slide,
    items: list[str],
    top=Inches(1.78),
    left=Inches(0.7),
    width=Inches(12.0),
    height=Inches(5.0),
    size=15,
) -> None:
    """Monospace commands with a light panel so curl/shell lines read clearly."""
    pad = Inches(0.06)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left - pad,
        top - pad,
        width + 2 * pad,
        height + 2 * pad,
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = C_CODE_BG
    panel.line.color.rgb = RGBColor(220, 225, 235)
    panel.line.width = Pt(0.75)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = FONT_CODE
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(35, 45, 55)
        p.font.bold = False
        try:
            p.line_spacing = 1.25
        except (TypeError, AttributeError):
            pass
        p.space_after = Pt(8)


def _section_slide(prs: Presentation, number: str, title: str, tagline: str) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, C_BG)
    _accent_bar(slide)
    # Large number
    nbox = slide.shapes.add_textbox(Inches(0.7), Inches(0.88), Inches(2.0), Inches(1.15))
    nt = nbox.text_frame
    nt.word_wrap = False
    np = nt.paragraphs[0]
    np.text = number
    np.font.name = FONT_TITLE
    np.font.size = Pt(68)
    np.font.bold = True
    np.font.color.rgb = C_TEAL
    tbox = slide.shapes.add_textbox(Inches(0.7), Inches(2.05), Inches(12.0), Inches(1.65))
    ttf = tbox.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.name = FONT_TITLE
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = C_NAVY
    try:
        tp.line_spacing = 1.15
    except (TypeError, AttributeError):
        pass
    sbox = slide.shapes.add_textbox(Inches(0.7), Inches(3.65), Inches(12.0), Inches(3.35))
    stf = sbox.text_frame
    stf.word_wrap = True
    stf.margin_left = Pt(4)
    stf.margin_right = Pt(4)
    # Split long taglines: sentences first, then comma clauses if still one wall of text
    parts = [p.strip() for p in re.split(r"(?<=[.!?])\s+", tagline) if p.strip()]
    if len(parts) == 1 and len(tagline) > 130:
        clauses = [c.strip() for c in tagline.split(",") if c.strip()]
        if len(clauses) >= 4:
            mid = max(2, len(clauses) // 2)
            parts = [", ".join(clauses[:mid]), ", ".join(clauses[mid:])]
        else:
            parts = [tagline]
    if not parts:
        parts = [tagline]
    for j, chunk in enumerate(parts):
        sp = stf.paragraphs[0] if j == 0 else stf.add_paragraph()
        sp.text = chunk
        sp.font.name = FONT_BODY
        sp.font.size = Pt(17)
        sp.font.color.rgb = C_MUTED
        sp.space_after = Pt(12)
        try:
            sp.line_spacing = 1.22
        except (TypeError, AttributeError):
            pass


def _card_row_title(slide, title: str) -> None:
    _title_block(slide, title)


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    out = root / "ppts" / "MLOps-2hr-Session-Consolidated-and-Execution.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Cover ----
    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _title_block(
        s,
        "MLOps Hands-On Workshop",
        "2 hours • Case studies + live mlops-workshop execution",
        top=Inches(1.1),
    )
    sub = s.shapes.add_textbox(Inches(0.7), Inches(2.85), Inches(12.0), Inches(1.0))
    sp0 = sub.text_frame.paragraphs[0]
    sp0.text = "Jenkins · Flask API · Financial ML pipeline · Automation"
    sp0.font.name = FONT_TITLE
    sp0.font.size = Pt(20)
    sp0.font.color.rgb = C_TEAL
    sp0.font.bold = True

    # ---- Session key points (exact pillars from brief) ----
    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Session Key Points (what we cover)")
    _bullets(
        s,
        [
            "Jenkins in MLOps — CI/CD and automated deployment",
            "Flask ML API for E-Commerce",
            "Cloud-Driven ML Pipeline for Financial Prediction",
            "Automated ML Pipelines and Workflow Automation",
        ],
        top=Inches(1.9),
        size=24,
    )
    foot = s.shapes.add_textbox(Inches(0.7), Inches(6.25), Inches(12.0), Inches(0.85))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "Flow: four pillars first, then hands-on with the mlops-workshop repo."
    fp.font.name = FONT_BODY
    fp.font.size = Pt(15)
    fp.font.color.rgb = C_MUTED
    fp.font.italic = True

    # ---- Agenda ----
    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Agenda (120 minutes)")
    _bullets(
        s,
        [
            "Part A — Four pillars (in order): Jenkins/CD; Flask e-commerce API; cloud financial ML; automated pipelines.",
            "Part B — Live lab: data → clean → train → container → API → GitHub Actions → Jenkins.",
            "Close — Security checks, versioning, rollback, monitoring.",
        ],
        top=Inches(1.85),
        size=20,
    )

    # ========== 1) JENKINS / MLOPS / CD ==========
    _section_slide(
        prs,
        "1",
        "Jenkins in MLOps — CI/CD and automated deployment",
        "Synthesized from: case study.pptx (Jenkins CD), Predictive Maintenance deck (Jenkins/GitHub in MLOps stack), "
        "and security/log-anomaly use case (where CD fits).",
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Why Jenkins (and CD) in MLOps?")
    _bullets(
        s,
        [
            "Continuous Delivery turns “model in a notebook” into repeatable stages: build → test → package → deploy.",
            "MLOps needs the same rigor as software: versioned artifacts, traceable promotions, fast rollback.",
            "Jenkins (or any CI) is the control plane: gates for tests, security checks, and deployment automation.",
            "In industry examples (e.g., predictive maintenance), stacks combine data versioning (e.g., DVC) with CI/CD (Jenkins, GitHub).",
        ],
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "What “good” looks like in the pipeline")
    _bullets(
        s,
        [
            "Single pipeline definition (Jenkinsfile / workflow YAML) = same steps every time.",
            "Stages: checkout → dependencies → data prep → train → test → build image → deploy → smoke test.",
            "Traceability: build number ↔ image tag ↔ model version ↔ deployment event.",
            "Security tie-in: CI/CD logs help detect anomalies (retries, failed gates, drift after deploy).",
        ],
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Bridge to today’s hands-on")
    _bullets(
        s,
        [
            "You will run the same lifecycle locally, then automate it with GitHub Actions and Jenkins.",
            "Goal: students see that ML delivery is “software delivery” with extra data and model artifacts.",
        ],
        size=22,
    )

    # ========== 2) FLASK E-COMMERCE API ==========
    _section_slide(
        prs,
        "2",
        "Flask ML API for E-Commerce",
        "Synthesized from: case study 3.pptx (Flask-enabled ML API for e-commerce recommendations).",
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Problem → solution (e-commerce)")
    _bullets(
        s,
        [
            "Platforms need real-time, scalable personalized recommendations (engagement, CTR, revenue).",
            "Pattern: train a model (e.g., collaborative filtering) and serve via a lightweight API.",
            "Flask provides REST routes such as recommendations per user and /health for operations.",
            "Model loaded in memory (or behind a worker) for low-latency responses; container + reverse proxy for scale.",
        ],
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "How this maps to mlops-workshop")
    _bullets(
        s,
        [
            "Our lab uses Flask with /predict, /health, /version — same serving pattern, different domain (fraud-style features).",
            "Emphasize: never log raw PII in request bodies; log metadata only (data security).",
            "Deployment: Docker image = portable “production slice” students can run on one machine.",
        ],
    )

    # ========== 3) CLOUD FINANCIAL ML PIPELINE ==========
    _section_slide(
        prs,
        "3",
        "Cloud-Driven ML Pipeline for Financial Prediction",
        "Synthesized from: case study (3).pptx (financial forecasting, quantization, cloud MLOps) + "
        "data-security / FinSecure-style fraud framing in data-security-usecase-demo-slides.pptx.",
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Financial ML: constraints and pipeline shape")
    _bullets(
        s,
        [
            "Financial prediction stresses cost, latency, and compliance: large models are expensive to run at scale.",
            "Cloud-driven pipelines standardize: ingest → features → train → evaluate → register → deploy → monitor.",
            "Optimization (e.g., quantization) reduces size and compute — relevant for batch and online inference.",
            "Data security: strict controls on PII, secrets, and what gets logged or stored in artifacts.",
        ],
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Tie-in: fraud / anomaly + governance")
    _bullets(
        s,
        [
            "Supervised scoring + operational monitoring (prediction distribution, errors) is standard in finance.",
            "CD reduces “works in dev, breaks in prod” and supports regulated change management.",
            "Today’s lab uses synthetic data with PII fields only to teach dropping PII before training.",
        ],
    )

    # ========== 4) AUTOMATED ML PIPELINES ==========
    _section_slide(
        prs,
        "4",
        "Automated ML Pipelines and Workflow Automation",
        "Synthesized from: case study (2).pptx (AutoML pipelines, Azure ML case study).",
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "What “automation” means in MLOps")
    _bullets(
        s,
        [
            "ML pipeline = repeatable workflow: ingest → preprocess → train → evaluate → deploy → monitor/retrain.",
            "AutoML accelerates model selection and tuning but still needs governance, testing, and deployment discipline.",
            "Benefits: speed, consistency, less manual error; tradeoffs: cost, opacity, cloud dependency.",
            "Tools named in source decks: Azure ML, SageMaker, Google AutoML, MLflow — concepts apply to Jenkins/GitHub too.",
        ],
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Workflow automation you will demonstrate")
    _bullets(
        s,
        [
            "GitHub Actions workflow: workflow-as-code, same stages as local shell scripts.",
            "Jenkins Pipeline: visual console, build history, integration with Docker for image + deploy.",
            "Both implement “automation” — choose by org standards; principles stay the same.",
        ],
    )

    # ========== PART B: EXECUTION (mlops-workshop) ==========
    _section_slide(
        prs,
        "B",
        "Hands-on: mlops-workshop execution",
        "Follow this slide deck while sharing terminal + browser; students use mlops-workshop-student clone.",
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Repository layout (what matters live)")
    _bullets(
        s,
        [
            "scripts/generate_synthetic_data.py — creates data/raw.csv (PII present on purpose for the lesson).",
            "src/data_prep.py — security cleaning → data/clean.csv",
            "src/train.py — artifacts/model.joblib, metrics, metadata",
            "src/app.py — Flask API; Dockerfile — image; Jenkinsfile + .github/workflows/mlops-ci.yml — automation",
        ],
        size=19,
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Execution — environment (start of demo block)")
    _command_block(
        s,
        [
            "cd mlops-workshop    # instructor (use mlops-workshop-student for class)",
            "python3 -m venv .venv && source .venv/bin/activate",
            "pip install -U pip && pip install -r requirements.txt",
            "chmod +x scripts/*.sh",
        ],
        top=Inches(1.72),
        size=14,
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Execution — full local pipeline (one shot)")
    _command_block(
        s,
        ["./scripts/pipeline_local_deploy.sh"],
        top=Inches(1.72),
        height=Inches(0.85),
        size=15,
    )
    _bullets(
        s,
        [
            "Or stepwise: generate_synthetic_data.py → data_prep.py → train.py → deploy.sh.",
            "Call out: PII dropped in data_prep.py before training.",
        ],
        top=Inches(2.95),
        height=Inches(3.9),
        size=19,
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Execution — verify the API (production slice)")
    _command_block(
        s,
        [
            "curl -s http://localhost:8000/health",
            "curl -s http://localhost:8000/version",
            'curl -s -X POST http://localhost:8000/predict \\',
            '  -H "Content-Type: application/json" \\',
            "  -d '{\"age\":42,\"country\":\"US\",\"device_type\":\"mobile\",...}'",
        ],
        top=Inches(1.65),
        height=Inches(2.35),
        size=13,
    )
    _bullets(
        s,
        [
            "Safe logging: app does not log raw request payloads.",
        ],
        top=Inches(4.35),
        height=Inches(2.5),
        size=19,
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Execution — GitHub Actions (automation #1)")
    _bullets(
        s,
        [
            "Open .github/workflows/mlops-ci.yml — walk stages (mirror Jenkins).",
            "Run: ./scripts/github_pipeline_sim.sh (local simulation of the workflow steps).",
            "Re-check /version after deploy to show a new build identity.",
        ],
        size=20,
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Execution — Jenkins (automation #2)")
    _bullets(
        s,
        [
            "cd jenkins && docker compose -f docker-compose.yml up --build -d",
            "Create Pipeline job from Jenkinsfile (mounted repo path as documented in jenkins/README.md).",
            "Build Now → console shows: generate → prep → train → test → docker build → deploy.",
        ],
        size=19,
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _card_row_title(s, "Execution — rollback & wrap (last 10–15 min)")
    _bullets(
        s,
        [
            "Optional: ./scripts/rollback.sh — redeploy previous image tag from deploy history.",
            "Recap: PII minimization, versioned artifacts, CI/CD as control plane, logs for anomaly awareness.",
            "Q&A + link to READING_LIST / student handout.",
        ],
        size=20,
    )

    s = _blank_slide(prs)
    _set_bg(s, C_BG)
    _accent_bar(s)
    _title_block(
        s,
        "Thank you",
        "Boring pipelines are good: predictable releases, safer ML in production.",
        top=Inches(2.35),
    )
    heart = s.shapes.add_textbox(Inches(0.7), Inches(4.05), Inches(12.0), Inches(1.2))
    hp = heart.text_frame.paragraphs[0]
    hp.text = "Deck file: MLOps-2hr-Session-Consolidated-and-Execution.pptx"
    hp.font.name = FONT_BODY
    hp.font.size = Pt(15)
    hp.font.color.rgb = C_MUTED

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
