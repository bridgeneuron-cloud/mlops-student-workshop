#!/usr/bin/env python3
"""
3-slide, visually rich summary of case study 3.pptx:
  Flask-Enabled ML API for E-Commerce Recommendations.

Source: mlops-workshop/ppts/case study 3.pptx (original multi-slide case study).

Requires: python-pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# E-commerce / product theme
C_BG = RGBColor(248, 250, 252)
C_HEADER = RGBColor(30, 58, 95)  # deep blue
C_ACCENT = RGBColor(232, 93, 4)  # warm orange
C_CARD_A = RGBColor(255, 255, 255)
C_CARD_BORDER = RGBColor(210, 220, 235)
C_MUTED = RGBColor(60, 65, 78)
C_TEAL = RGBColor(6, 122, 104)

FONT = "Calibri"


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _swidth(slide) -> int:
    return slide.part.package.presentation_part.presentation.slide_width


def _set_bg(slide, rgb: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def _bar(slide, color: RGBColor, h=Inches(0.2)) -> None:
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _swidth(slide), h
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def _tb(
    slide,
    left,
    top,
    width,
    height,
    title: str | None,
    lines: list[str],
    title_pt=22,
    body_pt=17,
    title_rgb=C_HEADER,
    body_rgb=C_MUTED,
    bold_title=True,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(title_pt)
        p.font.bold = bold_title
        p.font.color.rgb = title_rgb
        p.space_after = Pt(10)
        start = 0
    else:
        start = -1
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 and not title else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(body_pt)
        p.font.color.rgb = body_rgb
        p.space_after = Pt(8)
        try:
            p.line_spacing = 1.15
        except (TypeError, AttributeError):
            pass


def _card(
    slide,
    left,
    top,
    width,
    height,
    heading: str,
    bullets: list[str],
) -> None:
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = C_CARD_A
    r.line.color.rgb = C_CARD_BORDER
    r.line.width = Pt(1)
    _tb(
        slide,
        left + Inches(0.2),
        top + Inches(0.18),
        width - Inches(0.4),
        height - Inches(0.35),
        heading,
        bullets,
        title_pt=18,
        body_pt=15,
        title_rgb=C_HEADER,
    )


def _pill(slide, left, top, text: str, w=Inches(2.85)) -> None:
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.52))
    p.fill.solid()
    p.fill.fore_color.rgb = RGBColor(255, 243, 230)
    p.line.color.rgb = C_ACCENT
    p.line.width = Pt(0.5)
    tf = p.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.text = text
    para.font.name = FONT
    para.font.size = Pt(13)
    para.font.bold = True
    para.font.color.rgb = C_HEADER
    para.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    out = root / "ppts" / "Flask-E-Commerce-ML-API-3slides.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----- Slide 1: Hero title -----
    s = _blank(prs)
    _set_bg(s, C_BG)
    _bar(s, C_HEADER)
    accent = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.2), _swidth(s), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = C_ACCENT
    accent.line.fill.background()

    _tb(
        s,
        Inches(0.75),
        Inches(0.95),
        Inches(11.8),
        Inches(1.5),
        "Flask-Enabled ML API",
        [],
        title_pt=40,
        body_pt=17,
    )
    sub = s.shapes.add_textbox(Inches(0.75), Inches(1.85), Inches(11.8), Inches(0.9))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "E-Commerce Recommendations · Scalable personalized product suggestions"
    sp.font.name = FONT
    sp.font.size = Pt(20)
    sp.font.color.rgb = C_MUTED
    sp.font.italic = True

    _tb(
        s,
        Inches(0.75),
        Inches(2.95),
        Inches(11.8),
        Inches(1.2),
        "Case study focus",
        [
            "Deploy an ML model behind a lightweight REST API customers can call from your storefront.",
            "Same pattern as production MLOps: train → package → serve → monitor.",
        ],
        title_pt=20,
        body_pt=17,
    )

    # Three pills (visual chips)
    _pill(s, Inches(0.75), Inches(4.55), "Collaborative filtering / content-based ML")
    _pill(s, Inches(3.95), Inches(4.55), "Flask + JSON REST")
    _pill(s, Inches(7.15), Inches(4.55), "Docker · Gunicorn · Nginx")

    foot = s.shapes.add_textbox(Inches(0.75), Inches(6.35), Inches(11.8), Inches(0.6))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "Condensed from: case study 3.pptx"
    fp.font.name = FONT
    fp.font.size = Pt(12)
    fp.font.color.rgb = RGBColor(130, 135, 145)

    # ----- Slide 2: Challenge + solution (cards) -----
    s = _blank(prs)
    _set_bg(s, C_BG)
    _bar(s, C_HEADER)
    _tb(
        s,
        Inches(0.7),
        Inches(0.45),
        Inches(12.0),
        Inches(0.85),
        "Why it matters & what we build",
        [],
        title_pt=30,
    )
    _tb(
        s,
        Inches(0.7),
        Inches(1.15),
        Inches(12.0),
        Inches(1.35),
        None,
        [
            "E-commerce needs real-time, scalable personalization to lift engagement, CTR, and revenue.",
            "Gap: many stacks lack intelligent recommendations; we add ML + a serving layer the frontend can call.",
        ],
        title_pt=18,
        body_pt=16,
    )

    _card(
        s,
        Inches(0.7),
        Inches(2.55),
        Inches(3.75),
        Inches(2.35),
        "Data & features",
        [
            "User events: views, clicks, purchases",
            "Product metadata: category, price, brand",
            "Prep: missing values, scaling, encoding",
        ],
    )
    _card(
        s,
        Inches(4.65),
        Inches(2.55),
        Inches(3.95),
        Inches(2.35),
        "ML core",
        [
            "Matrix factorization or content-based model",
            "scikit-learn, pandas, numpy",
            "Metrics: RMSE, Precision@K · persist with joblib",
        ],
    )
    _card(
        s,
        Inches(8.85),
        Inches(2.55),
        Inches(3.75),
        Inches(2.35),
        "Flask API",
        [
            "GET /recommend/<user_id> → top-N items",
            "GET /health for ops",
            "Model in memory · JSON in/out",
        ],
    )

    strip = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(5.15),
        Inches(11.9),
        Inches(1.15),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(230, 245, 241)
    strip.line.color.rgb = C_TEAL
    strip.line.width = Pt(0.75)
    _tb(
        s,
        Inches(0.95),
        Inches(5.28),
        Inches(11.4),
        Inches(0.95),
        "Frontend integration",
        [
            "Async calls from the web or mobile app · personalized carousels · refresh as the shopper navigates.",
        ],
        title_pt=17,
        body_pt=15,
        title_rgb=C_TEAL,
    )

    # ----- Slide 3: Run in production -----
    s = _blank(prs)
    _set_bg(s, C_BG)
    _bar(s, C_HEADER)
    _tb(
        s,
        Inches(0.7),
        Inches(0.45),
        Inches(12.0),
        Inches(0.75),
        "Ship it: deploy, test, measure",
        [],
        title_pt=30,
    )

    # Left: stack diagram (boxes)
    y0 = Inches(1.2)
    step_h = Inches(0.72)
    gap = Inches(0.12)
    labels = [
        ("Docker image", "Portable runtime for dev/stage/prod"),
        ("Gunicorn", "WSGI workers for concurrent requests"),
        ("Nginx", "Reverse proxy & TLS termination"),
        ("Cloud", "EC2 · Heroku · Azure (your choice)"),
    ]
    for i, (t, subl) in enumerate(labels):
        top = y0 + i * (step_h + gap)
        bx = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.75),
            top,
            Inches(5.5),
            step_h,
        )
        bx.fill.solid()
        bx.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bx.line.color.rgb = C_CARD_BORDER
        _tb(s, Inches(0.95), top + Inches(0.1), Inches(5.1), Inches(0.55), t, [subl], title_pt=16, body_pt=12)

    flow = s.shapes.add_textbox(Inches(0.75), Inches(5.05), Inches(5.5), Inches(0.35))
    fp = flow.text_frame.paragraphs[0]
    fp.text = "Typical request path: Internet \u2192 Nginx \u2192 Gunicorn \u2192 Flask \u2192 model"
    fp.font.name = FONT
    fp.font.size = Pt(11)
    fp.font.color.rgb = C_TEAL
    fp.font.italic = True

    # Right: testing + outcomes
    _card(
        s,
        Inches(6.55),
        Inches(1.15),
        Inches(6.05),
        Inches(2.05),
        "Validate the API",
        [
            "Postman or curl · structured JSON request/response",
            "Logs for tracing (no sensitive payload dumps in prod)",
            "Optional: Prometheus metrics",
        ],
    )
    _card(
        s,
        Inches(6.55),
        Inches(3.35),
        Inches(6.05),
        Inches(2.05),
        "Outcomes & hard bits",
        [
            "Targets: higher engagement & CTR · sub-200 ms responses at scale",
            "Challenges: cold start, live data refresh, high concurrency",
            "Next: deep CF, Kafka streams, FastAPI",
        ],
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
