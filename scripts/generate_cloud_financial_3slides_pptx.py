#!/usr/bin/env python3
"""
3-slide, visually rich deck: Cloud-Driven ML Pipeline for Financial Prediction.

Sources synthesized:
  - Case Study_ Predictive Maintenance with MLOps.pptx (GE Aviation, MLOps stack, outcomes)
  - case study (3).pptx (cloud financial ML, quantization, pipeline architecture)
  - data-security-usecase-demo-slides.pptx (financial fraud / governance framing)

Requires: python-pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Financial / enterprise theme
C_BG = RGBColor(252, 250, 246)
C_HEADER = RGBColor(18, 52, 86)
C_ACCENT = RGBColor(180, 130, 16)  # gold
C_CARD_A = RGBColor(255, 255, 255)
C_CARD_BORDER = RGBColor(210, 215, 225)
C_MUTED = RGBColor(55, 58, 68)
C_TEAL = RGBColor(8, 105, 95)

FONT = "Calibri"


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _swidth(slide) -> int:
    return slide.part.package.presentation_part.presentation.slide_width


def _set_bg(slide, rgb: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def _bar(slide, color: RGBColor, h=Inches(0.2)) -> None:
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _swidth(slide), h
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def _tb(
    slide,
    left,
    top,
    width,
    height,
    title: str | None,
    lines: list[str],
    title_pt=22,
    body_pt=17,
    title_rgb=C_HEADER,
    body_rgb=C_MUTED,
    bold_title=True,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(title_pt)
        p.font.bold = bold_title
        p.font.color.rgb = title_rgb
        p.space_after = Pt(10)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 and not title else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(body_pt)
        p.font.color.rgb = body_rgb
        p.space_after = Pt(8)
        try:
            p.line_spacing = 1.15
        except (TypeError, AttributeError):
            pass


def _card(
    slide,
    left,
    top,
    width,
    height,
    heading: str,
    bullets: list[str],
    title_pt=16,
    body_pt=13,
) -> None:
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = C_CARD_A
    r.line.color.rgb = C_CARD_BORDER
    r.line.width = Pt(1)
    _tb(
        slide,
        left + Inches(0.18),
        top + Inches(0.16),
        width - Inches(0.36),
        height - Inches(0.32),
        heading,
        bullets,
        title_pt=title_pt,
        body_pt=body_pt,
        title_rgb=C_HEADER,
    )


def _pill(slide, left, top, text: str, w=Inches(3.9)) -> None:
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.5))
    p.fill.solid()
    p.fill.fore_color.rgb = RGBColor(255, 248, 230)
    p.line.color.rgb = C_ACCENT
    p.line.width = Pt(0.5)
    tf = p.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.text = text
    para.font.name = FONT
    para.font.size = Pt(12)
    para.font.bold = True
    para.font.color.rgb = C_HEADER
    para.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    out = root / "ppts" / "Cloud-Driven-ML-Pipeline-Financial-3slides.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----- Slide 1 -----
    s = _blank(prs)
    _set_bg(s, C_BG)
    _bar(s, C_HEADER)
    acc = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.2), _swidth(s), Inches(0.07)
    )
    acc.fill.solid()
    acc.fill.fore_color.rgb = C_ACCENT
    acc.line.fill.background()

    _tb(
        s,
        Inches(0.72),
        Inches(0.85),
        Inches(11.9),
        Inches(1.4),
        "Cloud-Driven ML Pipeline",
        [],
        title_pt=36,
    )
    sub = s.shapes.add_textbox(Inches(0.72), Inches(1.72), Inches(11.9), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Financial prediction · scale · cost-aware inference"
    sp.font.name = FONT
    sp.font.size = Pt(18)
    sp.font.color.rgb = C_MUTED
    sp.font.italic = True

    _tb(
        s,
        Inches(0.72),
        Inches(2.55),
        Inches(11.9),
        Inches(2.45),
        "Why “cloud-driven” for finance",
        [
            "Financial workloads need elastic compute, repeatable pipelines, and strong governance — cloud platforms deliver standardized ingest, training, and serving.",
            "Industrial MLOps case studies (e.g., predictive maintenance on sensor streams) show the same pattern: continuous data, versioned models, automated deploy, monitoring.",
            "For prediction use cases (risk, pricing, fraud), the pipeline must also control latency and cost — often via optimization and right-sized deployment.",
        ],
        title_pt=19,
        body_pt=15,
    )

    _pill(s, Inches(0.68), Inches(5.1), "Multi-cloud / SaaS ML")
    _pill(s, Inches(4.72), Inches(5.1), "CI/CD + versioned artifacts")
    _pill(s, Inches(8.75), Inches(5.1), "Secure serving & monitoring")

    foot = s.shapes.add_textbox(Inches(0.72), Inches(6.32), Inches(11.9), Inches(0.55))
    fp = foot.text_frame.paragraphs[0]
    fp.text = (
        "Sources: Case Study_ Predictive Maintenance with MLOps.pptx · case study (3).pptx · "
        "data-security-usecase-demo-slides.pptx"
    )
    fp.font.name = FONT
    fp.font.size = Pt(10)
    fp.font.color.rgb = RGBColor(125, 130, 140)

    # ----- Slide 2: Reference architectures -----
    s = _blank(prs)
    _set_bg(s, C_BG)
    _bar(s, C_HEADER)
    _tb(
        s,
        Inches(0.7),
        Inches(0.4),
        Inches(12.0),
        Inches(0.85),
        "Two lenses: industrial PdM and cloud financial ML",
        [],
        title_pt=27,
    )

    _card(
        s,
        Inches(0.68),
        Inches(1.22),
        Inches(3.95),
        Inches(2.55),
        "Predictive maintenance (reference)",
        [
            "Example: aviation — IoT sensor ingest (temp, vibration, pressure).",
            "MLOps stack (source deck): DVC · Jenkins / GitHub Actions · Kubernetes / Seldon · Prometheus / Grafana.",
            "Outcome story: proactive maintenance vs surprise failures.",
        ],
        title_pt=16,
        body_pt=13,
    )
    _card(
        s,
        Inches(4.78),
        Inches(1.22),
        Inches(3.95),
        Inches(2.55),
        "Cloud financial pipeline",
        [
            "Cloud ingest & prep → train/validate on VMs → deploy (e.g., serverless or managed endpoints).",
            "Monitoring & logging for drift, latency, and spend.",
            "Examples from financial ML deck: time-series / fraud-style workloads; PTQ/QAT to cut latency & cloud cost.",
        ],
        title_pt=16,
        body_pt=13,
    )
    _card(
        s,
        Inches(8.88),
        Inches(1.22),
        Inches(3.75),
        Inches(2.55),
        "Optimization angle",
        [
            "Quantization & distillation shrink models and speed inference.",
            "Illustrative tradeoff (source case): accuracy small dip vs large latency & cost win.",
            "Matches finance need: fast decisions on streaming or batch scoring.",
        ],
        title_pt=16,
        body_pt=13,
    )

    strip = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.68),
        Inches(3.95),
        Inches(11.95),
        Inches(1.2),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(232, 245, 242)
    strip.line.color.rgb = C_TEAL
    strip.line.width = Pt(0.75)
    _tb(
        s,
        Inches(0.9),
        Inches(4.05),
        Inches(11.5),
        Inches(1.0),
        "Financial services security framing",
        [
            "Real-time fraud / risk models demand data security: minimize PII in artifacts, safe logging, versioned deployments with rollback (FinSecure-style controls).",
        ],
        title_pt=16,
        body_pt=14,
        title_rgb=C_TEAL,
    )

    flow = s.shapes.add_textbox(Inches(0.7), Inches(5.35), Inches(12.0), Inches(0.45))
    fp = flow.text_frame.paragraphs[0]
    fp.text = "End-to-end arc: data → features → train → register → deploy → monitor → retrain (same story across PdM and finance)."
    fp.font.name = FONT
    fp.font.size = Pt(13)
    fp.font.color.rgb = C_MUTED
    fp.font.italic = True

    # ----- Slide 3: Impact + governance -----
    s = _blank(prs)
    _set_bg(s, C_BG)
    _bar(s, C_HEADER)
    _tb(
        s,
        Inches(0.7),
        Inches(0.4),
        Inches(12.0),
        Inches(0.75),
        "Proof points, trade-offs, and what to govern",
        [],
        title_pt=28,
    )

    _card(
        s,
        Inches(0.68),
        Inches(1.12),
        Inches(5.85),
        Inches(2.35),
        "Industrial MLOps outcomes (PdM deck)",
        [
            "Failure prediction accuracy up ~30% (reported).",
            "~$12M annual savings; maintenance downtime down ~40% (illustrative).",
            "Shows value of monitored, repeatable pipelines in production.",
        ],
        title_pt=16,
        body_pt=14,
    )
    _card(
        s,
        Inches(6.68),
        Inches(1.12),
        Inches(5.85),
        Inches(2.35),
        "Financial ML case (quantization deck)",
        [
            "Example metrics: latency improved materially; model size down; cloud cost down (~40% in cited scenario).",
            "Accuracy vs speed trade-off must be owned by the business and compliance.",
        ],
        title_pt=16,
        body_pt=14,
    )

    _card(
        s,
        Inches(0.68),
        Inches(3.62),
        Inches(11.85),
        Inches(1.65),
        "Governance checklist",
        [
            "Schema + PII controls before training · no secrets in logs · versioned model artifacts · auditable deploy/rollback.",
        ],
        title_pt=17,
        body_pt=15,
    )

    _tb(
        s,
        Inches(0.72),
        Inches(5.45),
        Inches(11.9),
        Inches(1.55),
        "Takeaway",
        [
            "Cloud-driven ML for finance borrows the same MLOps spine proven in industrial IoT: continuous delivery, observability, and cost-aware inference.",
            "Pair technical wins (latency, cost) with data security and regulatory-ready process.",
        ],
        title_pt=18,
        body_pt=15,
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
