#!/usr/bin/env python3
"""
3-slide, visually rich summary of case study (2).pptx:
  Automated ML Pipelines, Azure AutoML, telecom churn case study.

Note: The source deck is about AutoML / churn, not Flask e-commerce.
      Slide 3 includes a short bridge to REST/API serving (Flask or cloud endpoints).

Requires: python-pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Azure / automation feel
C_BG = RGBColor(248, 250, 252)
C_HEADER = RGBColor(36, 59, 83)
C_ACCENT = RGBColor(0, 114, 198)  # Azure-ish blue
C_CARD_A = RGBColor(255, 255, 255)
C_CARD_BORDER = RGBColor(210, 220, 235)
C_MUTED = RGBColor(60, 65, 78)
C_TEAL = RGBColor(6, 122, 104)

FONT = "Calibri"


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _swidth(slide) -> int:
    return slide.part.package.presentation_part.presentation.slide_width


def _set_bg(slide, rgb: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def _bar(slide, color: RGBColor, h=Inches(0.2)) -> None:
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _swidth(slide), h
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def _tb(
    slide,
    left,
    top,
    width,
    height,
    title: str | None,
    lines: list[str],
    title_pt=22,
    body_pt=17,
    title_rgb=C_HEADER,
    body_rgb=C_MUTED,
    bold_title=True,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(title_pt)
        p.font.bold = bold_title
        p.font.color.rgb = title_rgb
        p.space_after = Pt(10)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 and not title else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(body_pt)
        p.font.color.rgb = body_rgb
        p.space_after = Pt(8)
        try:
            p.line_spacing = 1.15
        except (TypeError, AttributeError):
            pass


def _card(
    slide,
    left,
    top,
    width,
    height,
    heading: str,
    bullets: list[str],
    title_pt=17,
    body_pt=14,
) -> None:
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = C_CARD_A
    r.line.color.rgb = C_CARD_BORDER
    r.line.width = Pt(1)
    _tb(
        slide,
        left + Inches(0.2),
        top + Inches(0.18),
        width - Inches(0.4),
        height - Inches(0.35),
        heading,
        bullets,
        title_pt=title_pt,
        body_pt=body_pt,
        title_rgb=C_HEADER,
    )


def _pill(slide, left, top, text: str, w=Inches(3.85)) -> None:
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.52))
    p.fill.solid()
    p.fill.fore_color.rgb = RGBColor(230, 242, 255)
    p.line.color.rgb = C_ACCENT
    p.line.width = Pt(0.5)
    tf = p.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.text = text
    para.font.name = FONT
    para.font.size = Pt(12)
    para.font.bold = True
    para.font.color.rgb = C_HEADER
    para.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    out = root / "ppts" / "Automated-ML-Pipelines-CaseStudy2-3slides.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----- Slide 1 -----
    s = _blank(prs)
    _set_bg(s, C_BG)
    _bar(s, C_HEADER)
    acc = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.2), _swidth(s), Inches(0.08)
    )
    acc.fill.solid()
    acc.fill.fore_color.rgb = C_ACCENT
    acc.line.fill.background()

    _tb(
        s,
        Inches(0.75),
        Inches(0.88),
        Inches(11.8),
        Inches(1.45),
        "Automated ML Pipelines",
        [],
        title_pt=38,
    )
    sub = s.shapes.add_textbox(Inches(0.75), Inches(1.78), Inches(11.8), Inches(0.85))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Revolutionizing data insights · MLOps · AutoML"
    sp.font.name = FONT
    sp.font.size = Pt(19)
    sp.font.color.rgb = C_MUTED
    sp.font.italic = True

    _tb(
        s,
        Inches(0.75),
        Inches(2.65),
        Inches(11.8),
        Inches(2.35),
        "Foundations",
        [
            "MLOps = Machine Learning + DevOps: automation, reproducibility, CI/CD for models.",
            "ML pipeline automates: ingest → preprocess → train → evaluate → deploy.",
            "AutoML pipelines add: model selection, hyperparameter tuning, feature engineering — faster experiments with less manual coding.",
        ],
        title_pt=20,
        body_pt=16,
    )

    _pill(s, Inches(0.7), Inches(5.15), "Automation & scale")
    _pill(s, Inches(4.75), Inches(5.15), "CI/CD for ML artifacts")
    _pill(s, Inches(8.8), Inches(5.15), "Deploy to an API endpoint")

    foot = s.shapes.add_textbox(Inches(0.75), Inches(6.35), Inches(11.8), Inches(0.55))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "Condensed from: case study (2).pptx  ·  Original deck: D. Naveen Reddy"
    fp.font.name = FONT
    fp.font.size = Pt(11)
    fp.font.color.rgb = RGBColor(130, 135, 145)

    # ----- Slide 2: Telecom churn case + workflow -----
    s = _blank(prs)
    _set_bg(s, C_BG)
    _bar(s, C_HEADER)
    _tb(
        s,
        Inches(0.7),
        Inches(0.42),
        Inches(12.0),
        Inches(1.0),
        "Case study: telecom churn with Azure ML + AutoML",
        [],
        title_pt=28,
    )
    _tb(
        s,
        Inches(0.7),
        Inches(1.2),
        Inches(12.0),
        Inches(0.95),
        None,
        [
            "Goal: predict which customers are likely to churn.",
            "Tooling: Azure Machine Learning and AutoML (source deck).",
        ],
        body_pt=16,
    )

    _card(
        s,
        Inches(0.7),
        Inches(2.35),
        Inches(3.75),
        Inches(2.45),
        "1 · Data",
        [
            "Customer profile, usage, complaints",
            "Clean and label: target column “Churn”",
        ],
    )
    _card(
        s,
        Inches(4.65),
        Inches(2.35),
        Inches(3.95),
        Inches(2.45),
        "2 · AutoML run",
        [
            "Multiple algorithms tried automatically",
            "Example outcome: XGBoost ~92% accuracy (per original case)",
        ],
    )
    _card(
        s,
        Inches(8.85),
        Inches(2.35),
        Inches(3.75),
        Inches(2.45),
        "3 · Ship",
        [
            "Evaluate and register the best model",
            "One-click style deploy to a scored API endpoint",
        ],
    )

    strip = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(5.05),
        Inches(11.9),
        Inches(1.25),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(230, 245, 241)
    strip.line.color.rgb = C_TEAL
    strip.line.width = Pt(0.75)
    _tb(
        s,
        Inches(0.95),
        Inches(5.18),
        Inches(11.4),
        Inches(1.05),
        "Same MLOps arc as a Flask lab",
        [
            "Train a model, then expose predictions over HTTP JSON — whether Azure-managed endpoint or your own Flask /predict route.",
        ],
        title_pt=17,
        body_pt=15,
        title_rgb=C_TEAL,
    )

    # ----- Slide 3: Benefits, challenges, tools -----
    s = _blank(prs)
    _set_bg(s, C_BG)
    _bar(s, C_HEADER)
    _tb(
        s,
        Inches(0.7),
        Inches(0.42),
        Inches(12.0),
        Inches(0.75),
        "Results, trade-offs, and tooling",
        [],
        title_pt=30,
    )

    _card(
        s,
        Inches(0.7),
        Inches(1.15),
        Inches(5.85),
        Inches(2.45),
        "Benefits (from the case study)",
        [
            "Faster model iteration · minimal custom code",
            "Easier retraining on new data",
            "Built-in experiment tracking and versioning culture",
        ],
    )
    _card(
        s,
        Inches(6.75),
        Inches(1.15),
        Inches(5.85),
        Inches(2.45),
        "Challenges",
        [
            "Large data can slow AutoML runs",
            "Less transparency into internal model choices",
            "Dependency on cloud pricing and quotas",
        ],
    )

    _card(
        s,
        Inches(0.7),
        Inches(3.75),
        Inches(11.9),
        Inches(1.55),
        "Representative tools",
        [
            "Azure ML · Google AutoML · Amazon SageMaker · MLflow (tracking) · Python (scikit-learn, pandas)",
        ],
        title_pt=17,
        body_pt=15,
    )

    _tb(
        s,
        Inches(0.75),
        Inches(5.45),
        Inches(11.8),
        Inches(1.65),
        "Takeaway",
        [
            "AutoML pipelines accelerate discovery; production still needs governance, tests, and a clear serving path.",
            "Link to Flask e-commerce APIs: both end with a model behind HTTP — choose managed cloud endpoints or your own Flask service for the same lesson.",
        ],
        title_pt=18,
        body_pt=15,
        title_rgb=C_HEADER,
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
