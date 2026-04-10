from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def add_title(slide, title: str, subtitle: str | None = None):
    # PPTX default title layout: title + subtitle
    title_shape = slide.shapes.title
    title_shape.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_bullets(slide, left, top, width, height, bullets: list[str], font_size=24):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size)
    return txBox


def add_header(slide, text: str):
    h = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.6))
    tf = h.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(28)
    return h


def main():
    prs = Presentation()

    # 1) Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(
        slide,
        "Real-time Data Security\nvia CI/CD + App Log Anomaly Detection",
        "Hands-on MLOps workshop • 2 hours • Local-only demo",
    )

    # 2) Why logs + why anomaly detection
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    add_header(slide, "Why logs? Why anomaly detection?")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(4.8),
        [
            "Security issues show up as behavioral change, not just model metrics.",
            "CI/CD logs reveal: build/deploy failures, dependency drift, unusual retries.",
            "App logs reveal: 5xx spikes, prediction errors, latency regressions.",
            "Anomaly detection turns “something feels off” into an actionable signal.",
        ],
    )

    # 3) Log sources
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_header(slide, "Log sources (what to monitor)")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(4.8),
        [
            "CI/CD: Jenkins/GitHub workflow stages, container build output, test summaries.",
            "Deployment: image tag, rollout/rollback events, health-check results.",
            "Application: /health, /version, /predict responses + error codes.",
            "Governance: who deployed what, and when (traceability).",
        ],
    )

    # 4) Example anomalies
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_header(slide, "Example anomalies (security-relevant)")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(4.8),
        [
            "Sudden increase in 401/403 or unexpected access patterns.",
            "Spike in 5xx or prediction schema validation errors after a deploy.",
            "CI pipeline retries beyond a normal baseline (dependency or secrets issue).",
            "Model version flips frequently (possible rollback/roll-forward loop).",
        ],
    )

    # 5) Data security: PII in logs
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_header(slide, "Data security rule: do not leak PII")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(4.8),
        [
            "Never log raw request payloads if they can contain PII.",
            "Log metadata instead (fields present, counts, schema shape).",
            "Keep artifacts versioned: you must explain inputs -> controls -> outputs.",
            "Apply the “security gate” before training and before logging.",
        ],
    )

    # 6) Detection architecture
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_header(slide, "Detection architecture (high-level)")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(4.8),
        [
            "Ingest logs from CI/CD, deployment, and app endpoints.",
            "Extract features: error rates, latency, retry counts, version transitions.",
            "Run anomaly detection: statistical baselines or ML-based detectors.",
            "Trigger: alert + automatic rollback or halt deploy pipeline.",
        ],
    )

    # 7) Where Jenkins/CD fits
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_header(slide, "Where Jenkins/CD fits")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(4.8),
        [
            "CI/CD is the “control plane”: repeatable steps + traceability.",
            "We gate security before artifacts are built and before deployment.",
            "We can test the API after deploy (health/version/predict).",
            "We can rollback by redeploying the previous image tag.",
        ],
    )

    # 8) What you will demo
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_header(slide, "What you will demo (local-only)")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(4.8),
        [
            "Synthetic dataset generation (includes PII to teach the security gate).",
            "Data security cleaning: drop PII columns before training.",
            "Train + save a versioned model artifact.",
            "Deploy behind a real-time API (/predict) and verify with curl.",
            "Run the same lifecycle via: GitHub copy + local Jenkins pipeline.",
        ],
        font_size=20,
    )

    # 9) Wrap
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_header(slide, "Wrap-up")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(4.8),
        [
            "Security is a pipeline property: control what you train on and what you log.",
            "CI/CD logs + app logs give the signals for anomaly detection.",
            "Versioned artifacts + rollback keeps deployments safe.",
            "Next step (optional): automate anomaly triggers for rollback/halt.",
        ],
    )

    out = "mlops-security-log-anomaly-usecase.pptx"
    prs.save(out)
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()

